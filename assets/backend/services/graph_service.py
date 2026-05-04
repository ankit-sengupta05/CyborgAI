"""
Graph Service — Knowledge graph with Leiden community detection.
Graphify-style visualization + Obsidian physics parameters.
"""
import asyncio
import hashlib
import structlog
from pathlib import Path
import networkx as nx

from services.embedding_service import EmbeddingService
from services.llm_service import LLMService
from services.database import GraphNodeDB, GraphEdgeDB, get_db

log = structlog.get_logger(__name__)


class GraphService:
    def __init__(self, embedding_service: EmbeddingService, llm_service: LLMService):
        self._embedding_svc = embedding_service
        self._llm_svc = llm_service
        self._graph = nx.Graph()
        self._nodes: dict[str, dict] = {}
        self._edges: list[dict] = []
        self._communities: dict[str, int] = {}
        self._community_meta: dict[int, dict] = {}

    async def initialize(self):
        """Load graph from database."""
        try:
            async with get_db() as db:
                nodes = await GraphNodeDB.get_all(db)
                edges = await GraphEdgeDB.get_all(db)

            for node in nodes:
                self._add_node_to_graph(node)
            for edge in edges:
                self._add_edge_to_graph(edge)

            if self._graph.number_of_nodes() > 0:
                await self._detect_communities()

            log.info(
                "Graph initialized",
                nodes=self._graph.number_of_nodes(),
                edges=self._graph.number_of_edges(),
            )
        except Exception as e:
            log.error(f"Graph initialization error: {e}")

    def _add_node_to_graph(self, node: dict):
        nid = node["id"]
        self._graph.add_node(nid, **node)
        self._nodes[nid] = node

    def _add_edge_to_graph(self, edge: dict):
        self._graph.add_edge(edge["source"], edge["target"], **edge)
        self._edges.append(edge)

    async def get_all(
        self, include_embeddings: bool = False, limit: int = 5000
    ) -> tuple[list[dict], list[dict], dict]:
        """Return all nodes and edges, optionally excluding embeddings for performance."""
        nodes_with_meta = []
        count = 0

        # Pre-calculate degrees for speed
        degrees = dict(self._graph.degree())

        for nid, data in self._nodes.items():
            if count >= limit:
                break
            community = self._communities.get(nid, 0)
            degree = degrees.get(nid, 0)

            node_data = {**data}
            if not include_embeddings:
                node_data.pop("embedding", None)

            # Truncate content for graph view to save memory
            if "content" in node_data and len(node_data["content"]) > 200:
                node_data["content"] = node_data["content"][:200] + "..."

            nodes_with_meta.append({
                **node_data,
                "community": community,
                "degree": degree,
            })
            count += 1

        # Limit edges too
        edges = self._edges[:limit*2]

        return nodes_with_meta, edges, self._community_meta

    async def search(self, query: str, limit: int = 20) -> list[dict]:
        """Semantic + keyword search over nodes."""
        if not self._nodes:
            return []

        results = []
        query_lower = query.lower()

        # Keyword search
        for nid, node in self._nodes.items():
            label_lower = node.get("label", "").lower()
            content_lower = node.get("content", "").lower()
            if query_lower in label_lower or query_lower in content_lower:
                results.append({**node, "score": 1.0})

        # Semantic search
        if self._embedding_svc.is_ready and len(results) < limit:
            query_emb = await self._embedding_svc.embed(query)
            candidate_embs = []
            candidate_ids = []
            for nid, node in self._nodes.items():
                emb = node.get("embedding")
                if emb:
                    candidate_embs.append(emb)
                    candidate_ids.append(nid)

            if candidate_embs:
                top_k = self._embedding_svc.top_k_similar(
                    query_emb, candidate_embs, k=limit
                )
                for idx, score in top_k:
                    if score > 0.5:
                        nid = candidate_ids[idx]
                        if nid not in {r["id"] for r in results}:
                            results.append({
                                **self._nodes[nid],
                                "score": score,
                            })

        return sorted(results, key=lambda x: x.get("score", 0), reverse=True)[:limit]

    async def ingest_file(self, file_path: str) -> dict:
        """Ingest a file into the knowledge graph."""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        if path.is_dir():
            return await self._ingest_directory(path)
        else:
            return await self._ingest_single_file(path)

    async def _ingest_directory(self, dir_path: Path) -> dict:
        total, success, failed = 0, 0, 0
        supported = {".txt", ".md", ".py", ".js", ".ts", ".json",
                     ".html", ".css", ".pdf", ".docx", ".pptx", ".csv"}
        for file_path in dir_path.rglob("*"):
            if file_path.is_file() and file_path.suffix.lower() in supported:
                total += 1
                try:
                    await self._ingest_single_file(file_path, detect_communities=False)
                    success += 1
                except Exception as e:
                    failed += 1
                    log.warning(f"Failed to ingest {file_path}: {e}")

        if success > 0:
            await self._detect_communities()

        return {"total": total, "success": success, "failed": failed}

    async def _ingest_single_file(self, file_path: Path, detect_communities: bool = True) -> dict:
        """Extract text, chunk, embed, and add to graph."""
        content = await self._extract_text(file_path)
        if not content.strip():
            return {"nodes_created": 0}

        chunks = self._chunk_text(content, chunk_size=800, overlap=100)
        nodes_created = 0

        for i, chunk in enumerate(chunks):
            # 1. Extract Triplets via LLM (Graphify style)
            triplets = await self._extract_triplets(chunk)

            # 2. Process extracted entities as nodes
            for entity in triplets.get("entities", []):
                # Deterministic ID based on normalized label for resolution
                label = entity.get("label", "Unknown").strip()
                node_id = hashlib.sha256(label.lower().encode()).hexdigest()[:16]

                if node_id not in self._nodes:
                    embedding = await self._embedding_svc.embed(chunk)
                    node = {
                        "id": node_id,
                        "label": label,
                        "content": chunk,
                        "content_type": entity.get(
                            "type", self._get_content_type(file_path.suffix)
                        ),
                        "source": str(file_path),
                        "chunk_index": i,
                        "embedding": embedding,
                        "community": 0,
                        "degree": 0,
                    }
                    async with get_db() as db:
                        await GraphNodeDB.create(db, node)
                    self._add_node_to_graph(node)
                    nodes_created += 1

            # 3. Process extracted relationships as edges
            for rel in triplets.get("relationships", []):
                s_label = rel.get("source", "").strip()
                t_label = rel.get("target", "").strip()
                if not s_label or not t_label:
                    continue

                s_id = hashlib.sha256(s_label.lower().encode()).hexdigest()[:16]
                t_id = hashlib.sha256(t_label.lower().encode()).hexdigest()[:16]

                if s_id in self._nodes and t_id in self._nodes:
                    edge = {
                        "source": s_id,
                        "target": t_id,
                        "type": rel.get("type", "related"),
                        "weight": rel.get("weight", 1.0),
                    }
                    async with get_db() as db:
                        await GraphEdgeDB.create(db, edge)
                    self._add_edge_to_graph(edge)

        # Re-detect communities after adding nodes

        # Re-detect communities after adding nodes
        if nodes_created > 0 and detect_communities:
            await self._detect_communities()

        return {"nodes_created": nodes_created, "file": str(file_path)}

    async def _extract_text(self, file_path: Path) -> str:
        """Extract text from various file types."""
        suffix = file_path.suffix.lower()
        loop = asyncio.get_event_loop()

        if suffix in (".txt", ".md", ".py", ".js", ".ts", ".json",
                      ".html", ".css", ".csv"):
            return file_path.read_text(encoding="utf-8", errors="ignore")

        elif suffix == ".pdf":
            return await loop.run_in_executor(None, self._extract_pdf, file_path)

        elif suffix == ".docx":
            return await loop.run_in_executor(None, self._extract_docx, file_path)

        elif suffix == ".pptx":
            return await loop.run_in_executor(None, self._extract_pptx, file_path)

        return ""

    def _extract_pdf(self, file_path: Path) -> str:
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(str(file_path))
            return "\n".join(
                p.extract_text() or "" for p in reader.pages
            )
        except Exception:
            return ""

    def _extract_docx(self, file_path: Path) -> str:
        try:
            from docx import Document
            doc = Document(str(file_path))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception:
            return ""

    def _extract_pptx(self, file_path: Path) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception:
            return ""

    def _chunk_text(self, text: str, chunk_size: int = 512,
                    overlap: int = 64) -> list[str]:
        """Split text into overlapping chunks."""
        words = text.split()
        if not words:
            return []
        chunks, start = [], 0
        while start < len(words):
            end = min(start + chunk_size, len(words))
            chunks.append(" ".join(words[start:end]))
            start += chunk_size - overlap
        return chunks

    def _extract_title(self, content: str, file_path: Path) -> str:
        lines = content.strip().split("\n")
        for line in lines:
            stripped = line.strip().lstrip("#").strip()
            if len(stripped) > 3:
                return stripped[:60]
        return file_path.stem[:60]

    def _get_content_type(self, suffix: str) -> str:
        mapping = {
            ".py": "code", ".js": "code", ".ts": "code",
            ".md": "markdown", ".pdf": "document",
            ".docx": "document", ".json": "data", ".csv": "data",
        }
        return mapping.get(suffix.lower(), "text")

    async def _extract_triplets(self, text: str) -> dict:
        """Use LLM to extract entities and relations in JSON format."""
        if not self._llm_svc.is_ready:
            return {"entities": [], "relationships": []}

        prompt = f"""
Extract key entities and their semantic relationships from the text below.
Focus on: (Entity A) -> [Relationship] -> (Entity B).
Keep entity labels short (1-3 words).

Use these specific relationship types:
- LINKED_TO (direct wikilink or reference)
- REFERENCED_BY (backlink or citation)
- SEMANTIC_SIMILAR (concepts in similar context)
- INFERRED (logically deduced connection)

Return ONLY valid JSON in this format:
{{
  "entities": [{{ "label": "Entity Name", "type": "person/org/concept/entity/file" }}],
  "relationships": [
    {{ "source": "Entity A", "target": "Entity B",
       "type": "LINKED_TO/INFERRED/...", "weight": 1.0 }}
  ]
}}

TEXT:
{text[:2000]}
"""
        try:
            response = await self._llm_svc.complete(
                prompt, temperature=0.1, max_tokens=1000
            )
            # Find JSON block
            import json
            import re
            match = re.search(r"\{.*\}", response, re.DOTALL)
            if match:
                return json.loads(match.group(0))
        except Exception as e:
            log.warning(f"LLM triplet extraction failed: {e}")

        return {"entities": [], "relationships": []}

    async def _detect_communities(self):
        """Leiden/Louvain community detection."""
        if self._graph.number_of_nodes() < 2:
            return

        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._run_community_detection)
        await self._analyze_communities()

    def _run_community_detection(self):
        try:
            import community as community_louvain
            partition = community_louvain.best_partition(self._graph)
            self._communities = partition
        except ImportError:
            for i, component in enumerate(
                nx.connected_components(self._graph)
            ):
                for node in component:
                    self._communities[node] = i

    async def _analyze_communities(self):
        """Generate semantic names and summaries for communities."""
        if not self._llm_svc.is_ready or not self._communities:
            return

        groups: dict[int, list[str]] = {}
        for nid, cid in self._communities.items():
            if cid not in groups:
                groups[cid] = []
            groups[cid].append(self._nodes[nid].get("label", ""))

        for cid, labels in groups.items():
            if cid in self._community_meta:
                continue

            context = ", ".join(labels[:20])
            prompt = f"Analyze these graph entities: {context}. " \
                     "Provide a 2-4 word 'name' and 1-sentence 'summary'. " \
                     "Return JSON: {\"name\": \"...\", \"summary\": \"...\"}"
            try:
                res = await self._llm_svc.complete(prompt, temperature=0.1)
                import json
                import re
                match = re.search(r"\{.*\}", res, re.DOTALL)
                if match:
                    self._community_meta[cid] = json.loads(match.group(0))
            except Exception:
                continue

    async def clear_graph(self):
        """Delete all nodes and edges from memory and database."""
        async with get_db() as db:
            await GraphNodeDB.delete_all(db)
        self._graph.clear()
        self._nodes.clear()
        self._edges.clear()
        self._communities.clear()
        log.info("Knowledge graph cleared")
