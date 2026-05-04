import asyncio
import structlog
from pathlib import Path
from typing import Optional
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

from config.settings import settings
from services.vault_service import VaultService
from services.voice_service import VoiceService

log = structlog.get_logger(__name__)


class InboxWatcher(FileSystemEventHandler):
    def __init__(self, service: 'IngestionService'):
        self.service = service

    def on_created(self, event):
        if event.is_directory:
            return
        log.info(f"New file in Inbox: {event.src_path}")
        asyncio.run_coroutine_threadsafe(
            self.service.process_file(event.src_path),
            self.service.loop
        )


class IngestionService:
    """
    Multimodal Ingestion Pipeline.
    Monitors Inbox/ and processes images, video, audio, and documents.
    """
    def __init__(self, vault_service: VaultService, voice_service: VoiceService):
        self.vault_service = vault_service
        self.voice_service = voice_service
        self.inbox_root = settings.brain_dir / "Inbox"
        self._is_running = False
        self._observer: Optional[Observer] = None
        self.loop = asyncio.get_event_loop()

    async def initialize(self):
        if self._is_running:
            return

        # Ensure subdirs exist
        for d in ["images", "videos", "audio", "documents"]:
            (self.inbox_root / d).mkdir(parents=True, exist_ok=True)

        self._is_running = True
        self._observer = Observer()
        self._observer.schedule(InboxWatcher(self), str(self.inbox_root), recursive=True)
        self._observer.start()
        log.info("Ingestion Service (Inbox Watcher) started", root=str(self.inbox_root))

    async def process_file(self, file_path: str):
        path = Path(file_path)
        if path.suffix.lower() == ".md":
            return  # Skip markdown

        log.info(f"Processing ingestion for {path.name}...")

        ext = path.suffix.lower()
        content = ""
        metadata = {}

        try:
            if ext in [".png", ".jpg", ".jpeg", ".webp"]:
                content = await self._process_image(path)
                metadata["type"] = "image"
            elif ext in [".mp4", ".mov", ".mkv", ".webm"]:
                content = await self._process_video(path)
                metadata["type"] = "video"
            elif ext in [".mp3", ".wav", ".m4a", ".ogg"]:
                content = await self._process_audio(path)
                metadata["type"] = "audio"
            elif ext in [".pdf", ".docx", ".pptx", ".txt"]:
                content = await self._process_document(path)
                metadata["type"] = "document"
            else:
                log.warning(f"Unsupported file type for ingestion: {ext}")
                return

            if content:
                # Create corresponding markdown note
                note_title = f"Ingested: {path.stem}"
                tags = ["ingested", metadata["type"]]

                # Organize into ACE structure or keep in Inbox?
                # PRD says ACE structure for permanent, but Inbox for staging.
                # Let's put in ACE/Atlas/Resources as default
                note = await self.vault_service.create_note(
                    title=note_title,
                    content=f"## Extracted Content from [[{path.name}]]\n\n{content}",
                    folder="atlas",  # ACE/Atlas
                    tags=tags,
                    note_type="ingestion",
                    area="Resources"
                )
                log.info(f"Ingestion complete: created note {note['id']}")

        except Exception as e:
            log.error(f"Ingestion failed for {path.name}: {e}")

    async def _process_image(self, path: Path) -> str:
        # TroCR or EasyOCR implementation placeholder
        # For now, let's use a placeholder message
        log.info(f"OCR processing (placeholder) for {path.name}")
        return f"[Image OCR Content Placeholder for {path.name}]"

    async def _process_audio(self, path: Path) -> str:
        # Use existing VoiceService (Whisper)
        if not self.voice_service.is_ready:
            return "Voice service not ready for transcription."

        # Whisper usually takes ndarray, but we can load from file
        import librosa
        audio, sr = librosa.load(str(path), sr=16000)
        transcript = self.voice_service.transcribe(audio)
        return transcript

    async def _process_video(self, path: Path) -> str:
        # Extract audio and transcribe
        log.info(f"Video processing (audio extraction) for {path.name}")
        return f"[Video Transcription Placeholder for {path.name}]"

    async def _process_document(self, path: Path) -> str:
        ext = path.suffix.lower()
        if ext == ".txt":
            return path.read_text(encoding="utf-8", errors="ignore")
        elif ext == ".docx":
            import docx
            doc = docx.Document(path)
            return "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        return f"[Document content for {path.name}]"

    async def cleanup(self):
        self._is_running = False
        if self._observer:
            self._observer.stop()
            self._observer.join()
