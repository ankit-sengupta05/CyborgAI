"""
Text extraction utility for document RAG / Multimodal chat.
Extracts text from PDF, DOCX, PPTX, TXT, CSV, MD.
"""
import os
import aiofiles
import structlog
from pathlib import Path
from typing import Optional

log = structlog.get_logger(__name__)


async def extract_text(file_path: str | Path) -> str:
    """Extract raw text from a document."""
    path = Path(file_path)
    if not path.exists():
        return ""

    ext = path.suffix.lower()
    
    try:
        if ext == ".pdf":
            return await _extract_pdf(path)
        elif ext == ".docx":
            return await _extract_docx(path)
        elif ext == ".pptx":
            return await _extract_pptx(path)
        elif ext in {".txt", ".md", ".csv"}:
            return await _extract_plain(path)
        else:
            log.warning(f"Unsupported document extension: {ext}")
            return ""
    except Exception as e:
        log.error(f"Error extracting text from {path.name}: {e}")
        return ""


async def _extract_pdf(path: Path) -> str:
    import PyPDF2
    text = []
    try:
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
    except Exception as e:
        log.warning(f"PyPDF2 failed: {e}")
    return "\n\n".join(text)


async def _extract_docx(path: Path) -> str:
    import docx
    try:
        doc = docx.Document(path)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        log.warning(f"python-docx failed: {e}")
        return ""


async def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    text = []
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
    except Exception as e:
        log.warning(f"python-pptx failed: {e}")
    return "\n".join(text)


async def _extract_plain(path: Path) -> str:
    """Read plain text files asynchronously."""
    async with aiofiles.open(path, 'r', encoding='utf-8', errors='ignore') as f:
        return await f.read()
