"""
Text Extractor Utility
Robust multi-format extraction with fallback chains.

PDF strategy:
  1. PyMuPDF (fitz)  — fastest, best text + layout
  2. pdfplumber      — good for tables / complex layouts
  3. PyPDF2          — legacy fallback
  4. pytesseract OCR — last resort for scanned/image-only pages

Image OCR strategy:
  1. pytesseract     — fast, no internet, good for printed text
  2. TrOCR           — handwriting (lazy-loaded, falls back if not installed)
"""
import asyncio
import structlog
from pathlib import Path

log = structlog.get_logger(__name__)

# ─── Lazy singletons for heavy models ─────────────────────────────────────────
_trocr_processor = None
_trocr_model = None


def _load_trocr():
    global _trocr_processor, _trocr_model
    if _trocr_processor is None:
        try:
            from transformers import TrOCRProcessor, VisionEncoderDecoderModel
            log.info("Loading TrOCR model (first use)…")
            _trocr_processor = TrOCRProcessor.from_pretrained(
                "microsoft/trocr-base-printed"  # lighter than handwritten
            )
            _trocr_model = VisionEncoderDecoderModel.from_pretrained(
                "microsoft/trocr-base-printed"
            )
            log.info("TrOCR model loaded ✓")
        except Exception as e:
            log.warning(f"TrOCR unavailable: {e}")
    return _trocr_processor, _trocr_model


# ─── Public API ───────────────────────────────────────────────────────────────

async def extract_text(file_path: Path) -> str:
    """Async entry point — runs blocking extractors in executor."""
    ext = file_path.suffix.lower()
    loop = asyncio.get_event_loop()

    if ext in (".txt", ".md", ".py", ".js", ".ts", ".json",
               ".html", ".css", ".csv"):
        try:
            return file_path.read_text(encoding="utf-8", errors="ignore")
        except Exception as e:
            log.error(f"Text read failed for {file_path.name}: {e}")
            return ""

    if ext == ".pdf":
        return await loop.run_in_executor(None, _extract_pdf, file_path)

    if ext == ".docx":
        return await loop.run_in_executor(None, _extract_docx, file_path)

    if ext == ".pptx":
        return await loop.run_in_executor(None, _extract_pptx, file_path)

    if ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif"):
        return await loop.run_in_executor(None, _extract_image, file_path)

    log.debug(f"Unsupported extension for extraction: {ext}")
    return ""


# ─── PDF extraction (fallback chain) ─────────────────────────────────────────

def _extract_pdf(file_path: Path) -> str:
    # 1. Try PyMuPDF (fitz) — best quality
    text = _pdf_pymupdf(file_path)
    if text and len(text.strip()) > 50:
        return text

    # 2. Try pdfplumber
    text = _pdf_pdfplumber(file_path)
    if text and len(text.strip()) > 50:
        return text

    # 3. Try PyPDF2
    text = _pdf_pypdf2(file_path)
    if text and len(text.strip()) > 50:
        return text

    # 4. OCR fallback for scanned / image-only PDFs
    log.info(f"PDF text extraction yielded little content for {file_path.name} — trying OCR")
    return _pdf_ocr_fallback(file_path)


def _pdf_pymupdf(file_path: Path) -> str:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(file_path))
        pages = []
        for page in doc:
            pages.append(page.get_text("text"))
        doc.close()
        return "\n".join(pages)
    except ImportError:
        pass
    except Exception as e:
        log.debug(f"PyMuPDF failed for {file_path.name}: {e}")
    return ""


def _pdf_pdfplumber(file_path: Path) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(str(file_path)) as pdf:
            pages = []
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    pages.append(t)
            return "\n".join(pages)
    except ImportError:
        pass
    except Exception as e:
        log.debug(f"pdfplumber failed for {file_path.name}: {e}")
    return ""


def _pdf_pypdf2(file_path: Path) -> str:
    try:
        import PyPDF2
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            return "\n".join(p.extract_text() or "" for p in reader.pages)
    except ImportError:
        pass
    except Exception as e:
        log.debug(f"PyPDF2 failed for {file_path.name}: {e}")
    return ""


def _pdf_ocr_fallback(file_path: Path) -> str:
    """Convert each PDF page to image and OCR it."""
    pages_text = []
    try:
        # Try pdf2image → pytesseract
        from pdf2image import convert_from_path
        images = convert_from_path(str(file_path), dpi=200)
        for i, img in enumerate(images):
            text = _ocr_image_pil(img)
            if text.strip():
                pages_text.append(f"[Page {i+1}]\n{text}")
        if pages_text:
            return "\n\n".join(pages_text)
    except ImportError:
        log.debug("pdf2image not installed — cannot OCR PDF pages")
    except Exception as e:
        log.warning(f"PDF OCR failed for {file_path.name}: {e}")

    # Last resort: return filename as placeholder so graph still creates a node
    return f"[Scanned PDF: {file_path.name} — install pdf2image and tesseract for OCR]"


# ─── Image extraction ─────────────────────────────────────────────────────────

def _extract_image(file_path: Path) -> str:
    """OCR an image file. Try pytesseract first, then TrOCR."""
    # 1. pytesseract (fast, good for printed text)
    text = _ocr_tesseract_path(file_path)
    if text and len(text.strip()) > 10:
        return f"Extracted Text (OCR):\n{text}"

    # 2. TrOCR (better for handwriting, heavier)
    text = _ocr_trocr(file_path)
    if text and len(text.strip()) > 10:
        return f"Extracted Text (TrOCR):\n{text}"

    return f"[Image OCR: no text detected in {file_path.name}]"


def _ocr_image_pil(pil_image) -> str:
    """OCR a PIL image object."""
    try:
        import pytesseract
        return pytesseract.image_to_string(pil_image)
    except ImportError:
        pass
    except Exception as e:
        log.debug(f"pytesseract (PIL) failed: {e}")
    return ""


def _ocr_tesseract_path(file_path: Path) -> str:
    """OCR a file path directly with pytesseract."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(str(file_path)).convert("RGB")
        return pytesseract.image_to_string(img)
    except ImportError:
        log.debug("pytesseract not installed")
    except Exception as e:
        log.debug(f"pytesseract failed for {file_path.name}: {e}")
    return ""


def _ocr_trocr(file_path: Path) -> str:
    """Run TrOCR on an image file."""
    try:
        from PIL import Image
        processor, model = _load_trocr()
        if processor is None or model is None:
            return ""

        image = Image.open(str(file_path)).convert("RGB")
        pixel_values = processor(images=image, return_tensors="pt").pixel_values
        generated_ids = model.generate(pixel_values, max_new_tokens=256)
        return processor.batch_decode(generated_ids, skip_special_tokens=True)[0]
    except Exception as e:
        log.debug(f"TrOCR failed for {file_path.name}: {e}")
    return ""


# ─── Document extractors ──────────────────────────────────────────────────────

def _extract_docx(file_path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(file_path))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                parts.append("\t".join(c.text.strip() for c in row.cells if c.text.strip()))
        return "\n".join(parts)
    except ImportError:
        log.error("python-docx not installed")
        return f"[Error: python-docx not installed to read {file_path.name}]"
    except Exception as e:
        log.error(f"DOCX extraction failed for {file_path.name}: {e}")
        return ""


def _extract_pptx(file_path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
                if hasattr(shape, "has_table") and shape.has_table:
                    for row in shape.table.rows:
                        parts.append(
                            "\t".join(c.text.strip() for c in row.cells if c.text.strip())
                        )
        return "\n".join(parts)
    except ImportError:
        log.error("python-pptx not installed")
        return f"[Error: python-pptx not installed to read {file_path.name}]"
    except Exception as e:
        log.warning(f"PPTX extraction failed for {file_path.name}: {e}")
        return ""
